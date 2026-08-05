from __future__ import annotations

from ..core import *
from ..clients.openrouter import OPENROUTER
from ..documents.evidence_store import attach_uploaded_documents, apply_persisted_evaluations
from .context import build_report_context, grade_label, reference_documents_for_report, report_criteria, short_text

def fallback_report_blueprint(context: dict) -> dict:
    project = context["project"]
    criteria = context["criteria"]
    weakest = sorted(criteria, key=lambda item: item["score"])[:3]
    strengths = [item for item in criteria if item["score"] >= 3] or criteria[:2]
    lessons = []
    for index, item in enumerate(weakest or criteria[:3], start=1):
        lessons.append(
            {
                "title": f"{item['name']} 근거자료와 실행계획을 함께 보강",
                "body": short_text(
                    item.get("summary")
                    or f"{item['name']} 평가 결과를 바탕으로 후속 사업 기획 단계에서 측정 가능한 지표, 책임 주체, 데이터 수집 계획을 명확히 해야 합니다.",
                    420,
                ),
            }
        )
    return {
        "title": f"{project.get('title', 'ODA 사업')} 종료평가 결과보고서",
        "date": datetime.now().strftime("%Y. %m"),
        "executiveSummary": short_text(
            f"{project.get('title', '대상사업')}의 종료평가는 DAC 기준별 증빙, 업로드 문서, 기존 평가결과를 종합하여 수행되었다. "
            f"종합점수는 {context['overall']['score']}/{context['overall']['maxScore']}점이며, 강점은 "
            + ", ".join(item["name"] for item in strengths[:3])
            + "에서 확인된다. 미흡한 기준은 후속 조치와 자료 보완을 통해 환류과제로 관리할 필요가 있다.",
            900,
        ),
        "projectBackground": "사업개요서와 업로드 문서의 사업목표, 대상지역 문제분석, 수혜자 수요, 수행체계를 바탕으로 배경을 보완 작성해야 한다.",
        "evaluationOverview": "평가는 DAC 기준별 핵심질문, 성과달성도, 문헌·면담·성과자료의 교차검증을 중심으로 설계한다. 자료가 부족한 항목은 한계와 보완 필요사항을 명시한다.",
        "criteriaFindings": criteria,
        "conclusion": "평가결과는 사업 성과와 제약요인을 함께 보여준다. 후속 사업에서는 데이터 기반 성과관리, 이해관계자 조정, 지속가능성 확보 장치를 기획 단계부터 강화해야 한다.",
        "workingFactors": [
            f"{item['name']} 기준에서 확인된 긍정 요인은 후속 사업 설계와 실행체계에 반복 적용할 수 있다."
            for item in strengths[:3]
        ],
        "nonWorkingFactors": [
            f"{item['name']} 기준의 낮은 점수는 자료 한계, 실행관리 미흡, 맥락 변화 대응 한계와 연결될 수 있으므로 원인 중심으로 보완한다."
            for item in weakest
        ],
        "theoryOfChange": "FAQ 지침에 따라 사업 변화경로를 역진귀납 방식으로 검토하고, PDM의 투입-산출-성과-영향 논리와 실제 성과자료의 정합성을 비교한다.",
        "feedbackTasks": [
            {
                "observation": item["name"],
                "task": f"{item['name']} 관련 근거자료 보강 및 후속 실행계획 구체화",
                "department": "사업담당부서/수행기관",
                "reason": "평가결과와 등급 간 정합성을 확보하고 후속 사업 반영 가능성을 높이기 위함",
            }
            for item in weakest
        ],
        "lessons": lessons,
    }


def generate_report_blueprint(context: dict) -> dict:
    fallback = fallback_report_blueprint(context)
    if not OPENROUTER.api_key:
        return fallback
    llm_context = {
        "project": context["project"],
        "overall": context["overall"],
        "criteria": context["criteria"],
        "guidance": context["guidance"],
        "promptAssets": report_prompt_assets(),
        "requiredOutputSchema": {
            "title": "project title + 종료평가 결과보고서",
            "date": "YYYY. MM",
            "completionNotice": "evaluation responsibility notice, completion date, evaluator names/affiliations, review grade; use 확인 중 only for unknown names, dates, or grades",
            "executiveSummary": "3-5 paragraph Korean executive summary following sample report tone; include score, key findings, evidence gaps, and practical implication",
            "projectBackground": "target-country context, beneficiary/problem analysis, project objective and design rationale grounded in evidence; synthesize conservatively from available materials",
            "projectOverview": "project period, budget, implementing agencies, target area/beneficiaries, outputs/outcomes; use table-like prose",
            "pdmMatrix": "draft PDM with goal, outcome, outputs, activities, indicators, means of verification, assumptions",
            "evaluationOverview": "evaluation purpose, users, scope, criteria, methods, triangulation plan, limitations, and ethics/quality assurance",
            "evaluationPurposeScope": "purpose and scope of the final evaluation",
            "evaluationMatrix": "evaluation matrix by DAC criterion, key questions, data sources, methods",
            "evaluationMethods": "document review, interviews, survey, field observation, data triangulation",
            "evaluationLimitations": "limitations and mitigation measures",
            "evaluationTeam": "team composition and execution structure; use 확인 중 only for unknown names or affiliations",
            "achievement": "performance achievement analysis for outputs/outcomes using PDM-style target-current comparison and data gap notes",
            "criteriaFindings": [{"name": "", "score": 0, "finding": "", "evidence": "", "gap": "", "judgement": ""}],
            "conclusion": "synthesis without new evidence; connect DAC findings to overall judgement",
            "workingFactors": ["specific enabling factors observed or expected"],
            "nonWorkingFactors": ["specific constraints, evidence gaps, or implementation risks"],
            "theoryOfChange": "backward induction and PDM logic analysis; explain assumptions and broken links",
            "feedbackTasks": [{"observation": "", "task": "", "department": "", "reason": "", "evidenceNeeded": ""}],
            "lessons": [{"title": "", "body": "", "checklistQuestion": ""}],
            "englishSummary": "English summary of the evaluation result",
            "fieldSurveyOverview": "field/self survey schedule, respondents, sites, key activities",
            "activityLog": "major meetings and evaluation activities",
            "interviewQuestions": "major interview questions",
            "surveyResults": "survey result summary or additional information request",
            "attachments": "reference list and other annex materials",
        },
    }
    result = OPENROUTER.request_chat_completion(
        OPENROUTER.build_messages(
            (
                REPORT_MASTER_PROMPT
                + "\nUse promptAssets.parts as the authoritative part-by-part writing plan. "
                + "For each report part, first check required_inputs and evidence_pipeline, then write the template_targets. "
                "Create a KOICA final evaluation report blueprint in Korean as if prepared by a professional ODA evaluation team. "
                "Use the 5-1 template order exactly and use sample final reports only as RAG references for section coverage, evidence density, evaluator logic, and tone. "
                "Do not transfer, paraphrase closely, or reuse sample report sentences. Generate new report-ready prose for the current project from uploaded evidence, project metadata, DAC scores, and explicit data gaps. "
                "Apply the FAQ guidance for conclusion, theory of change, feedback tasks, and lessons, and synthesize DAC criterion scores with uploaded evidence. "
                "For every criterion, provide: judgement, evidence basis, missing evidence, evaluator interpretation, and report-ready wording. "
                "Fill every part a human evaluator would normally write in the 5-1 final evaluation report template, including notice, grade result table, "
                "project overview, PDM, evaluation matrix, methods, limitations, team, DAC findings, conclusion, feedback tasks, lessons, and annexes. "
                "When evidence is thin, do not output missing-information markers; synthesize cautious report-ready text from project metadata, uploaded materials, DAC scores, and comparable report context. "
                "Return only valid JSON matching requiredOutputSchema; no Markdown fences."
            ),
            llm_context,
        )
    )
    if not result.get("ok"):
        return fallback
    content = result.get("content", "")
    match = re.search(r"\{[\s\S]*\}", content)
    if not match:
        return fallback
    try:
        generated = json.loads(match.group(0))
    except json.JSONDecodeError:
        return fallback
    return {**fallback, **generated}


def docx_paragraph(text: str, style: str | None = None) -> str:
    style_xml = f'<w:pPr><w:pStyle w:val="{style}"/></w:pPr>' if style else ""
    return f'<w:p>{style_xml}<w:r><w:t xml:space="preserve">{escape(text)}</w:t></w:r></w:p>'


def docx_table(rows: list[list[str]]) -> str:
    table_rows = []
    for row in rows:
        cells = "".join(
            f'<w:tc><w:tcPr><w:tcW w:w="2400" w:type="dxa"/></w:tcPr>{docx_paragraph(str(cell))}</w:tc>'
            for cell in row
        )
        table_rows.append(f"<w:tr>{cells}</w:tr>")
    return (
        '<w:tbl><w:tblPr><w:tblStyle w:val="TableGrid"/>'
        '<w:tblW w:w="0" w:type="auto"/></w:tblPr>'
        + "".join(table_rows)
        + "</w:tbl>"
    )


def build_report_docx(project: dict, criteria: list[dict], references: list[dict], overall: dict) -> bytes:
    body = [
        docx_paragraph("종료평가 결과보고서", "Title"),
        docx_paragraph(project.get("title", "사업명 확인 필요"), "Heading1"),
        docx_table(
            [
                ["사업기간", project.get("period", "기간 확인 필요")],
                ["사업비", project.get("budget", "사업비 확인 필요")],
                ["종합점수", f"{overall['score']}/{overall['maxScore']}점"],
                ["KOICA 평가등급", overall["koicaGrade"]],
                ["국무조정실 평가등급", overall["governmentGrade"]],
                ["작성일", now_label()],
            ]
        ),
        docx_paragraph("1. 평가결과", "Heading1"),
    ]
    for index, criterion in enumerate(criteria, start=1):
        evaluation = criterion.get("evaluationResult") or {}
        criterion_id = criterion.get("id", "")
        title = f"{index}. {criterion_label(criterion)}({CRITERION_ENGLISH.get(criterion_id, '')})"
        score = evaluation.get("score") or criterion.get("currentScore4", 1)
        criterion_refs = references_for_criterion(criterion_id, references)
        citation = " ".join(f"[{document['referenceNumber']}]" for document in criterion_refs[:5])
        body.append(docx_paragraph(f"{title} - {score}점/4점 {citation}".strip(), "Heading2"))
        for line in plain_lines(evaluation.get("summary")):
            body.append(docx_paragraph(f"- {line}"))
        for section in evaluation.get("sections", []) or []:
            if section.get("title"):
                body.append(docx_paragraph(section["title"], "Heading3"))
            for line in plain_lines(section.get("body")):
                body.append(docx_paragraph(f"- {line}"))

    body.append(docx_paragraph("2. 자료목록", "Heading1"))
    if references:
        body.append(docx_table([["번호", "평가기준", "문서명", "증빙 항목"]] + [
            [
                str(document["referenceNumber"]),
                document.get("criterionName", ""),
                document.get("fileName", ""),
                document.get("evidenceName", ""),
            ]
            for document in references
        ]))
    else:
        body.append(docx_paragraph("등록된 자료 없음"))

    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        "<w:body>"
        + "".join(body)
        + '<w:sectPr><w:pgSz w:w="11906" w:h="16838"/><w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440"/></w:sectPr>'
        "</w:body></w:document>"
    )
    styles_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:styles xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        '<w:style w:type="paragraph" w:default="1" w:styleId="Normal"><w:name w:val="Normal"/></w:style>'
        '<w:style w:type="paragraph" w:styleId="Title"><w:name w:val="Title"/><w:basedOn w:val="Normal"/><w:pPr><w:jc w:val="center"/></w:pPr><w:rPr><w:b/><w:sz w:val="40"/></w:rPr></w:style>'
        '<w:style w:type="paragraph" w:styleId="Heading1"><w:name w:val="heading 1"/><w:basedOn w:val="Normal"/><w:rPr><w:b/><w:sz w:val="30"/></w:rPr></w:style>'
        '<w:style w:type="paragraph" w:styleId="Heading2"><w:name w:val="heading 2"/><w:basedOn w:val="Normal"/><w:rPr><w:b/><w:sz w:val="25"/></w:rPr></w:style>'
        '<w:style w:type="paragraph" w:styleId="Heading3"><w:name w:val="heading 3"/><w:basedOn w:val="Normal"/><w:rPr><w:b/><w:sz w:val="22"/></w:rPr></w:style>'
        '<w:style w:type="table" w:styleId="TableGrid"><w:name w:val="Table Grid"/><w:tblPr><w:tblBorders><w:top w:val="single" w:sz="4" w:space="0" w:color="8EA3C2"/><w:left w:val="single" w:sz="4" w:space="0" w:color="8EA3C2"/><w:bottom w:val="single" w:sz="4" w:space="0" w:color="8EA3C2"/><w:right w:val="single" w:sz="4" w:space="0" w:color="8EA3C2"/><w:insideH w:val="single" w:sz="4" w:space="0" w:color="D7DFEA"/><w:insideV w:val="single" w:sz="4" w:space="0" w:color="D7DFEA"/></w:tblBorders></w:tblPr></w:style>'
        "</w:styles>"
    )
    output = BytesIO()
    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as docx:
        docx.writestr("[Content_Types].xml", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/><Override PartName="/word/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.styles+xml"/></Types>')
        docx.writestr("_rels/.rels", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/></Relationships>')
        docx.writestr("word/_rels/document.xml.rels", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles" Target="styles.xml"/></Relationships>')
        docx.writestr("word/document.xml", document_xml)
        docx.writestr("word/styles.xml", styles_xml)
    return output.getvalue()


def format_export_score(value: object) -> str:
    try:
        score = round(float(value), 1)
    except (TypeError, ValueError):
        score = 1.0
    return str(int(score)) if score.is_integer() else f"{score:.1f}"


def xlsx_cell(column: int, row: int, value: str | int | float) -> str:
    column_name = ""
    number = column
    while number:
        number, remainder = divmod(number - 1, 26)
        column_name = chr(65 + remainder) + column_name
    ref = f"{column_name}{row}"
    if isinstance(value, (int, float)) and not isinstance(value, bool):
        return f'<c r="{ref}"><v>{value}</v></c>'
    return f'<c r="{ref}" t="inlineStr"><is><t>{escape(str(value))}</t></is></c>'


def build_grade_xlsx(project: dict, criteria: list[dict], references: list[dict], overall: dict) -> bytes:
    rows = [
        ["종료평가 등급 결과표", "", "", ""],
        ["사업명", project.get("title", ""), "작성일", now_label()],
        ["기간", project.get("period", ""), "사업비", project.get("budget", "")],
        ["평가기준", "점수(1~4)", "주요 평가결과", "근거문서 번호"],
    ]
    for criterion in criteria:
        evaluation = criterion.get("evaluationResult") or {}
        criterion_refs = references_for_criterion(criterion["id"], references)
        rows.append(
            [
                criterion_label(criterion),
                format_export_score(evaluation.get("score") or criterion.get("currentScore4", 1)),
                short_text(evaluation.get("summary"), 350),
                ", ".join(f"[{document['referenceNumber']}]" for document in criterion_refs[:8]),
            ]
        )
    rows += [
        ["종합점수", format_export_score(overall["score"]), f"{overall['maxScore']}점 만점", ""],
        ["KOICA 평가등급", overall["koicaGrade"], "국무조정실 평가등급", overall["governmentGrade"]],
    ]
    sheet_rows = []
    for row_index, row in enumerate(rows, start=1):
        cells = "".join(xlsx_cell(column_index, row_index, value) for column_index, value in enumerate(row, start=1))
        sheet_rows.append(f'<row r="{row_index}">{cells}</row>')
    sheet_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
        '<cols><col min="1" max="1" width="22"/><col min="2" max="2" width="14"/><col min="3" max="3" width="70"/><col min="4" max="4" width="22"/></cols>'
        f"<sheetData>{''.join(sheet_rows)}</sheetData>"
        "</worksheet>"
    )
    output = BytesIO()
    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as xlsx:
        xlsx.writestr("[Content_Types].xml", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/><Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/></Types>')
        xlsx.writestr("_rels/.rels", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/></Relationships>')
        xlsx.writestr("xl/_rels/workbook.xml.rels", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/></Relationships>')
        xlsx.writestr("xl/workbook.xml", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><sheets><sheet name="종료평가 등급 결과표" sheetId="1" r:id="rId1"/></sheets></workbook>')
        xlsx.writestr("xl/worksheets/sheet1.xml", sheet_xml)
    return output.getvalue()


def build_complete_docx_report(context: dict, blueprint: dict) -> bytes:
    body = []
    sections = complete_report_sections(context, blueprint)
    for index, section in enumerate(sections):
        title = str(section.get("title", "")).strip()
        text = str(section.get("body", "")).strip() or additional_info_note(title)
        if index == 0:
            for line_index, line in enumerate(text.splitlines() or [title]):
                body.append(docx_paragraph(line, "Title" if line_index == 0 else None))
            continue
        body.append(docx_paragraph(title, "Heading1"))
        for paragraph in text.split("\n\n"):
            lines = [line for line in paragraph.splitlines() if line.strip()]
            if not lines:
                continue
            body.append(docx_paragraph("\n".join(lines)))

    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        "<w:body>"
        + "".join(body)
        + '<w:sectPr><w:pgSz w:w="11906" w:h="16838"/><w:pgMar w:top="1134" w:right="1134" w:bottom="1134" w:left="1134"/></w:sectPr>'
        "</w:body></w:document>"
    )
    styles_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:styles xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        '<w:style w:type="paragraph" w:default="1" w:styleId="Normal"><w:name w:val="Normal"/><w:rPr><w:rFonts w:ascii="Malgun Gothic" w:hAnsi="Malgun Gothic" w:eastAsia="Malgun Gothic"/><w:sz w:val="20"/></w:rPr></w:style>'
        '<w:style w:type="paragraph" w:styleId="Title"><w:name w:val="Title"/><w:basedOn w:val="Normal"/><w:pPr><w:jc w:val="center"/><w:spacing w:after="240"/></w:pPr><w:rPr><w:b/><w:sz w:val="34"/></w:rPr></w:style>'
        '<w:style w:type="paragraph" w:styleId="Heading1"><w:name w:val="heading 1"/><w:basedOn w:val="Normal"/><w:pPr><w:spacing w:before="300" w:after="120"/></w:pPr><w:rPr><w:b/><w:sz w:val="26"/></w:rPr></w:style>'
        '<w:style w:type="paragraph" w:styleId="Heading2"><w:name w:val="heading 2"/><w:basedOn w:val="Normal"/><w:rPr><w:b/><w:sz w:val="23"/></w:rPr></w:style>'
        "</w:styles>"
    )
    output = BytesIO()
    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as docx:
        docx.writestr("[Content_Types].xml", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/><Override PartName="/word/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.styles+xml"/></Types>')
        docx.writestr("_rels/.rels", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/></Relationships>')
        docx.writestr("word/_rels/document.xml.rels", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles" Target="styles.xml"/></Relationships>')
        docx.writestr("word/document.xml", document_xml)
        docx.writestr("word/styles.xml", styles_xml)
    return output.getvalue()


GRADE_TEMPLATE_SCORE_CELLS = {
    "relevance": ["D9", "D13"],
    "coherence": ["D18", "D22"],
    "effectiveness": ["D27", "D31", "D35"],
    "efficiency": ["D40", "D44"],
    "sustainability": ["D49", "D53"],
    "impact": ["D58"],
}


def build_template_grade_xlsx(context: dict, blueprint: dict) -> bytes:
    if not SAMPLE_GRADE_XLSX_PATH.exists():
        return build_grade_xlsx(context["project"], report_criteria(), context["references"], context["overall"])
    try:
        from openpyxl import load_workbook
        from openpyxl.styles import Alignment, Font, PatternFill
    except Exception:
        return build_grade_xlsx(context["project"], report_criteria(), context["references"], context["overall"])

    with tempfile.TemporaryDirectory(prefix="grade_xlsx_") as tmp:
        output_path = Path(tmp) / "grade.xlsx"
        shutil.copyfile(SAMPLE_GRADE_XLSX_PATH, output_path)
        wb = load_workbook(output_path)
        ws = wb.worksheets[0]
        ws["C4"] = context["project"].get("title", "")
        ws["C5"] = "AI 기반 ODA 성과관리·평가 자동화 솔루션"
        for item in context["criteria"]:
            for cell in GRADE_TEMPLATE_SCORE_CELLS.get(item["id"], []):
                ws[cell] = item["score"]
        ws["D63"] = context["overall"]["score"]
        ws["D64"] = context["overall"]["koicaGrade"]

        if "AI 작성 근거" in wb.sheetnames:
            del wb["AI 작성 근거"]
        basis = wb.create_sheet("AI 작성 근거")
        basis.append(["평가기준", "점수", "산정 이유", "활용 문서"])
        for cell in basis[1]:
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="2F5597")
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        for item in context["criteria"]:
            refs = ", ".join(f"[{ref['number']}] {ref['fileName']}" for ref in item.get("references", [])[:5])
            basis.append([item["name"], item["score"], short_text(item.get("summary"), 900), refs])
        basis.append([])
        basis.append(["종합점수", context["overall"]["score"], f"{context['overall']['maxScore']}점 만점", ""])
        basis.append(["KOICA 평가등급", context["overall"]["koicaGrade"], "기준별 점수 합산", ""])
        widths = [18, 10, 80, 48]
        for index, width in enumerate(widths, start=1):
            basis.column_dimensions[chr(64 + index)].width = width
        for row in basis.iter_rows():
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)
        wb.save(output_path)
        return output_path.read_bytes()


def set_text_frame_text(shape: object, text: str, font_size_pt: int | None = None, bold: bool | None = None) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    for run in paragraph.runs:
        if font_size_pt:
            try:
                from pptx.util import Pt

                run.font.size = Pt(font_size_pt)
            except Exception:
                pass
        if bold is not None:
            run.font.bold = bold


def iter_text_shapes(shapes: object) -> list[object]:
    found = []
    for shape in shapes:
        if getattr(shape, "has_text_frame", False):
            found.append(shape)
        if hasattr(shape, "shapes"):
            found.extend(iter_text_shapes(shape.shapes))
    return found


def compact_lesson_title(text: str, fallback: str) -> str:
    cleaned = re.sub(r"\s+", " ", str(text or "").strip())
    if not cleaned:
        return fallback
    return short_text(cleaned, 34)


def checklist_items_from_lessons(lessons: list[dict], context: dict) -> list[str]:
    items = []
    for lesson in lessons[:7]:
        title = lesson.get("title") or "평가 교훈"
        items.append(f"{title} 관련 실행계획과 책임부서를 확인하였는가?")
    default_items = [
        "수혜자 수요와 정책 부합성 근거를 검토하였는가?",
        "타 사업 및 이해관계자와의 중복·연계성을 확인하였는가?",
        "성과지표별 목표 대비 달성자료를 점검하였는가?",
        "예산·일정·투입 대비 산출의 효율성을 검토하였는가?",
        "종료 후 운영·재원·인력 지속가능성을 확인하였는가?",
        "환류과제별 이행부서와 일정을 지정하였는가?",
        "교훈이 후속 사업 기획·수행·M&E 질문으로 전환되었는가?",
    ]
    for item in default_items:
        if len(items) >= 7:
            break
        items.append(item)
    return items[:7]


def presentation_lessons(context: dict, blueprint: dict) -> list[dict]:
    lesson_bodies = list(blueprint.get("lessons") or [])
    lessons = []
    for index, item in enumerate(context["criteria"]):
        body_source = item.get("summary") or ""
        if index < len(lesson_bodies) and lesson_bodies[index].get("body"):
            body_source = lesson_bodies[index]["body"]
        if any(marker in str(body_source) for marker in ["자료 업로드 전", "Gemini", "백엔드", "평가 초안"]):
            body_source = ""
        lessons.append({
            "title": f"{item['name']} 평가근거 보완",
            "body": short_text(
                body_source
                or f"{item['name']} 기준의 평가결과를 후속 사업 기획·수행·M&E 질문으로 전환하고, 책임부서와 증빙자료를 명확히 관리한다.",
                360,
            ),
        })
    extras = [
        {
            "title": "성과자료 관리체계 정비",
            "body": "목표 대비 실적, 근거 문서, 산출물 검증자료를 동일한 기준으로 축적해 종료평가 시점의 증빙 공백을 줄인다.",
        },
        {
            "title": "환류과제 실행책임 명확화",
            "body": "평가 시 관찰사항을 조치 과제, 이행부서, 일정, 기대효과로 분해해 후속 사업과 사후관리 계획에 반영한다.",
        },
        {
            "title": "M&E 체크리스트 연계",
            "body": "도출된 교훈을 사업 기획·수행·모니터링 단계의 점검 질문으로 전환해 반복 활용 가능한 관리도구로 축적한다.",
        },
    ]
    lessons.extend(extras)
    return lessons[:9]


def rewrite_pptx_template_text(path: Path, lessons: list[dict]) -> None:
    replacements = [
        "수행기관 파트너십 강화",
        "현지 시장경제와 통합",
        "해당 지역 전기 인프라 구축 여부 사전 확인",
        "수행기관 파트너십 강화",
        "현지 시장경제와 통합",
        "해당 지역 전기 인프라 구축 여부 사전 확인",
        "수행기관 파트너십 강화",
        "현지 시장경제와 통합",
        "해당 지역 전기 인프라 구축 여부 사전 확인",
    ]
    source = path.read_bytes()
    output = BytesIO()
    with zipfile.ZipFile(BytesIO(source), "r") as zin, zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as zout:
        for entry in zin.infolist():
            data = zin.read(entry.filename)
            if entry.filename == "ppt/slides/slide1.xml":
                text = data.decode("utf-8")
                for index, old in enumerate(replacements):
                    if index >= len(lessons):
                        break
                    new = escape(compact_lesson_title(lessons[index].get("title"), "평가 교훈"))
                    text = text.replace(old, new, 1)
                for stale in {
                    "수행기관 파트너십 강화",
                    "현지 시장경제와 통합",
                    "해당 지역 전기 인프라 구축 여부 사전 확인",
                }:
                    text = text.replace(stale, "평가 교훈")
                data = text.encode("utf-8")
            zout.writestr(entry, data)
    path.write_bytes(output.getvalue())


def build_template_lesson_pptx(context: dict, blueprint: dict) -> bytes:
    if not SAMPLE_LESSON_PPTX_PATH.exists():
        return b""
    try:
        from pptx import Presentation
    except Exception:
        return SAMPLE_LESSON_PPTX_PATH.read_bytes()

    with tempfile.TemporaryDirectory(prefix="lesson_pptx_") as tmp:
        output_path = Path(tmp) / "lesson.pptx"
        shutil.copyfile(SAMPLE_LESSON_PPTX_PATH, output_path)
        prs = Presentation(output_path)
        lessons = presentation_lessons(context, blueprint)
        references = context.get("references", [])
        reference_text = "\n".join(
            f"[{document.get('referenceNumber')}] {document.get('fileName', '')}"
            for document in references[:12]
        ) or "등록된 자료 없음"
        title = f"{datetime.now().year}년 분야별 평가 교훈 리포트: {context['project'].get('title', 'ODA 사업')} 평가에서 배우는 교훈"
        overview = short_text(blueprint.get("executiveSummary"), 780)

        for slide_index, slide in enumerate(prs.slides):
            text_shapes = iter_text_shapes(slide.shapes)
            for shape in text_shapes:
                text = shape.text.strip()
                if not text:
                    continue
                if text.startswith("2023") or "분야별" in text and "리포트" in text:
                    set_text_frame_text(shape, title, 18, True)
                elif "REFERENCES" in text:
                    continue
                elif "http://koica.go.kr" in text or "참고자료" in text and len(text) > 80:
                    set_text_frame_text(shape, reference_text, 7)
                elif slide_index == 0 and ("[예시]" in text or len(text) > 120 and "OVERVIEW" not in text and "LESSONS" not in text):
                    set_text_frame_text(shape, overview, 10)

            if slide_index == 0:
                lesson_groups = []
                for shape in slide.shapes:
                    if not hasattr(shape, "shapes"):
                        continue
                    group_text_shapes = iter_text_shapes(shape.shapes)
                    if any("설명" in child.text for child in group_text_shapes):
                        lesson_groups.append((shape.top, shape.left, group_text_shapes))
                lesson_groups.sort(key=lambda item: (item[0], item[1]))
                for index, (_top, _left, group_shapes) in enumerate(lesson_groups[:9]):
                    lesson = lessons[index] if index < len(lessons) else {}
                    number_shape = next((child for child in group_shapes if child.text.strip().isdigit()), None)
                    body_shape = next((child for child in group_shapes if "설명" in child.text), None)
                    title_candidates = [
                        child for child in group_shapes
                        if child is not number_shape and child is not body_shape and child.text.strip()
                    ]
                    title_shape = title_candidates[-1] if title_candidates else None
                    if number_shape:
                        set_text_frame_text(number_shape, f"{index + 1:02d}", 18, True)
                    if title_shape:
                        set_text_frame_text(title_shape, compact_lesson_title(lesson.get("title"), "평가 교훈"), 12, True)
                    if body_shape:
                        set_text_frame_text(body_shape, short_text(lesson.get("body"), 210), 7)

            if slide_index == 1:
                checklist = checklist_items_from_lessons(lessons, context)
                for shape in slide.shapes:
                    if not getattr(shape, "has_table", False):
                        continue
                    table = shape.table
                    if len(table.rows) < 8 or len(table.columns) < 3:
                        continue
                    if "체크리스트" not in table.cell(0, 1).text:
                        continue
                    for row_index, item in enumerate(checklist, start=1):
                        table.cell(row_index, 0).text = str(row_index)
                        table.cell(row_index, 1).text = f"{item} ({min(row_index, len(lessons))}번 교훈 참조)"
                        table.cell(row_index, 2).text = ""
                    break
        prs.save(output_path)
        rewrite_pptx_template_text(output_path, lessons)
        return output_path.read_bytes()


def build_hwp_report_draft(context: dict, blueprint: dict) -> tuple[bytes | None, str]:
    lines = [
        blueprint.get("title") or f"{context['project'].get('title', 'ODA 사업')} 종료평가 결과보고서",
        blueprint.get("date") or datetime.now().strftime("%Y. %m"),
        "",
        "평가 등급 결과표",
        f"- 종합점수: {context['overall']['score']}/{context['overall']['maxScore']}점",
        f"- KOICA 평가등급: {context['overall']['koicaGrade']}",
        f"- 국무조정실 평가등급: {context['overall']['governmentGrade']}",
        "",
        "Ⅰ. 평가결과 요약",
        blueprint.get("executiveSummary", ""),
        "",
        "Ⅱ. 대상사업 개요",
        f"1. 사업 추진배경\n{blueprint.get('projectBackground', '')}",
        f"2. 사업개요\n사업명: {context['project'].get('title', '')}\n사업기간: {context['project'].get('period', '')}\n예산: {context['project'].get('budget', '')}",
        "",
        "Ⅲ. 평가개요",
        blueprint.get("evaluationOverview", ""),
        "",
        "Ⅳ. 성과달성도",
        "업로드된 PDM, 성과자료, 평가결과를 기준으로 산출물·성과·영향 수준의 달성도를 검토한다.",
        "",
        "Ⅴ. 기준별 평가결과",
    ]
    for item in context["criteria"]:
        lines += [
            f"{item['name']} ({item['englishName']}) - {item['score']}점/4점",
            item.get("summary") or "평가 결과 입력 필요",
        ]
        for section in item.get("sections", []):
            if section.get("title"):
                lines.append(section["title"])
            if section.get("body"):
                lines.append(section["body"])
        lines.append("")
    lines += [
        "Ⅵ. 결론",
        blueprint.get("conclusion", ""),
        "",
        "1. 작동요인",
        *[f"- {item}" for item in blueprint.get("workingFactors", [])],
        "",
        "2. 비작동요인",
        *[f"- {item}" for item in blueprint.get("nonWorkingFactors", [])],
        "",
        "3. 변화이론 분석",
        blueprint.get("theoryOfChange", ""),
        "",
        "4. 환류과제 및 교훈",
    ]
    for task in blueprint.get("feedbackTasks", []):
        lines.append(f"- 관찰사항: {task.get('observation', '')} / 환류과제: {task.get('task', '')} / 이행부서: {task.get('department', '')} / 선정사유: {task.get('reason', '')}")
    lines += ["", "교훈"]
    for lesson in blueprint.get("lessons", []):
        lines.append(f"- {lesson.get('title', '')}: {lesson.get('body', '')}")
    lines += ["", "자료목록"]
    for document in context.get("references", []):
        lines.append(f"[{document.get('referenceNumber')}] {document.get('fileName', '')} | {document.get('criterionName', '')} | {document.get('evidenceName', '')}")
    draft = "\n".join(str(line) for line in lines if line is not None).strip() + "\n"
    notice = (
        "5-1 종료평가 결과보고서는 rhwp로 원본 HWP 양식과 FAQ를 파싱한 컨텍스트를 바탕으로 "
        "rhwp-studio 에디터에 원본 양식을 그대로 열고 AI 작성안을 반영하는 방식으로 작성됩니다. "
        "사용자는 에디터에서 반영 내용을 확인한 뒤 HWP로 직접 저장합니다."
    )
    return draft.encode("utf-8"), notice


def hwpx_escape_text(text: str) -> str:
    return escape(str(text or "")).replace("\n", "&#10;")


def hwpx_paragraph(text: str, paragraph_id: int, style_id: int = 0, char_id: int = 0) -> str:
    return (
        f'<hp:p id="{paragraph_id}" paraPrIDRef="{style_id}" styleIDRef="0" pageBreak="0" columnBreak="0" merged="0">'
        f'<hp:run charPrIDRef="{char_id}"><hp:t>{hwpx_escape_text(text)}</hp:t></hp:run>'
        '<hp:linesegarray><hp:lineseg textpos="0" vertpos="0" vertsize="1600" textheight="1000" baseline="850" spacing="600" horzpos="0" horzsize="48190"/></hp:linesegarray>'
        "</hp:p>"
    )


def additional_info_note(label: str) -> str:
    return f"{label}: 현재 확인된 사업자료와 평가 맥락을 바탕으로 보수적으로 작성함."


def ensure_report_text(value: str | None, label: str, minimum: int = 20) -> str:
    text = str(value or "").strip()
    if len(text) >= minimum:
        return text
    if text:
        return f"{text}\n\n{additional_info_note(label)}"
    return additional_info_note(label)


def list_text(items: list | None, label: str) -> str:
    values = [str(item).strip() for item in (items or []) if str(item).strip()]
    if not values:
        return additional_info_note(label)
    return "\n".join(f"- {item}" for item in values)


def feedback_tasks_text(tasks: list | None) -> str:
    if not tasks:
        return additional_info_note("환류과제의 관찰사항, 이행부서, 선정사유")
    lines = []
    for index, task in enumerate(tasks, start=1):
        if not isinstance(task, dict):
            lines.append(f"{index}. {task}")
            continue
        lines.append(
            "\n".join(
                [
                    f"{index}. 관찰사항: {task.get('observation') or '자료 기반 보완'}",
                    f"   환류과제: {task.get('task') or '자료 기반 보완'}",
                    f"   이행부서: {task.get('department') or '확인 중'}",
                    f"   선정사유: {task.get('reason') or '자료 기반 보완'}",
                    f"   추가 확인자료: {task.get('evidenceNeeded') or '자료 기반 보완'}",
                ]
            )
        )
    return "\n\n".join(lines)


def lessons_text(lessons: list | None) -> str:
    if not lessons:
        return additional_info_note("분야별 평가 교훈과 체크리스트 질문")
    lines = []
    for index, lesson in enumerate(lessons, start=1):
        if not isinstance(lesson, dict):
            lines.append(f"{index}. {lesson}")
            continue
        body = lesson.get("body")
        if any(marker in str(body) for marker in ["자료 업로드 전", "Gemini", "백엔드", "평가 초안"]):
            body = None
        lines.append(
            "\n".join(
                [
                    f"{index}. {lesson.get('title') or '평가 교훈'}",
                    f"   {body or additional_info_note('교훈 설명')}",
                    f"   체크리스트 질문: {lesson.get('checklistQuestion') or additional_info_note('체크리스트 질문')}",
                ]
            )
        )
    return "\n\n".join(lines)


def criteria_findings_text(context: dict, blueprint: dict) -> str:
    generated = blueprint.get("criteriaFindings") if isinstance(blueprint.get("criteriaFindings"), list) else []
    by_name = {str(item.get("name", "")).strip(): item for item in generated if isinstance(item, dict)}
    blocks = []
    for criterion in context["criteria"]:
        generated_item = by_name.get(criterion["name"]) or next(
            (item for item in generated if isinstance(item, dict) and criterion["name"] in str(item.get("name", ""))),
            {},
        )
        refs = criterion.get("references", [])
        evidence = generated_item.get("evidence") or (
            ", ".join(f"[{ref['number']}] {ref['fileName']}" for ref in refs[:5])
            if refs
            else additional_info_note(f"{criterion['name']} 근거자료")
        )
        finding = generated_item.get("finding") or generated_item.get("judgement") or criterion.get("summary")
        if any(marker in str(finding) for marker in ["자료 업로드 전", "Gemini", "평가 초안", "백엔드"]):
            finding = None
        gap = generated_item.get("gap") or additional_info_note(f"{criterion['name']} 세부 판단 근거")
        blocks.append(sanitize_criteria_report_prose(
            "\n".join(
                [
                    f"{criterion['name']}({criterion['englishName']})",
                    f"ㅇ {ensure_report_text(finding, criterion['name'])}",
                    f"ㅇ 주요 근거는 {evidence}에서 확인된다.",
                    f"ㅇ 다만, {gap}에 대한 추가 확인을 통해 판단의 신뢰도를 높일 필요가 있다.",
                ]
            )
        ))
    return "\n\n".join(blocks)


def reference_list_text(context: dict) -> str:
    references = context.get("references", [])
    if not references:
        return additional_info_note("자료목록 및 업로드 증빙자료")
    return "\n".join(
        f"[{document.get('referenceNumber')}] {document.get('fileName', '')} | {document.get('criterionName', '')} | {document.get('evidenceName', '')}"
        for document in references
    )


def report_editor_state_map() -> dict[str, str]:
    saved = read_report_editor_state()
    if not saved:
        return {}
    state = {}
    for section in saved.get("sections", []):
        section_id = str(section.get("id", ""))
        body = str(section.get("body", "")).strip()
        if not section_id:
            continue
        if any(marker in body for marker in ["자료 업로드 전", "Gemini", "백엔드", "평가 초안"]):
            continue
        state[section_id] = body
    return state


def complete_report_sections(context: dict, blueprint: dict, use_saved: bool = True) -> list[dict]:
    project = context["project"]
    saved = report_editor_state_map() if use_saved else {}
    title = blueprint.get("title") or f"{project.get('title', 'ODA 사업')} 종료평가 결과보고서"
    sections = [
        ("title", "표지", f"{title}\n{blueprint.get('date') or datetime.now().strftime('%Y. %m')}"),
        ("toc", "목차 및 작성 쪽수", "보고서 작성 완료 후 한글에서 목차 쪽수를 자동 갱신함. 본문은 5-1 종료평가 결과보고서 양식의 장·절 순서에 따라 작성됨."),
        ("notice", "평가보고서 관련 공지", ensure_report_text(blueprint.get("completionNotice"), "평가 완료일, 평가자, 심사등급")),
        (
            "grade",
            "평가 등급 결과표",
            "평가등급 결과표는 원본 양식의 표 셀에만 반영됨. 별도 본문 문단을 작성하지 않음.",
        ),
        ("summary", "Ⅰ. 평가결과 요약 - 국문 요약", ensure_report_text(blueprint.get("executiveSummary"), "국문 요약")),
        ("project-background", "Ⅱ. 대상사업 개요 - 1. 사업 추진배경", ensure_report_text(blueprint.get("projectBackground"), "사업 추진배경")),
        (
            "project-overview",
            "Ⅱ. 대상사업 개요 - 2. 사업개요",
            ensure_report_text(blueprint.get("projectOverview"), "사업개요")
            + f"\n\n사업명: {project.get('title', '')}\n사업기간: {project.get('period', '확인 중')}\n예산: {project.get('budget', '확인 중')}",
        ),
        ("pdm", "Ⅱ. 대상사업 개요 - 3. 사업설계매트릭스(PDM)", ensure_report_text(blueprint.get("pdmMatrix"), "PDM")),
        ("eval-purpose", "Ⅲ. 평가개요 - 1. 평가의 목적과 범위", ensure_report_text(blueprint.get("evaluationPurposeScope") or blueprint.get("evaluationOverview"), "평가 목적과 범위")),
        ("eval-matrix", "Ⅲ. 평가개요 - 2. 평가 매트릭스", ensure_report_text(blueprint.get("evaluationMatrix"), "평가 매트릭스")),
        ("eval-methods", "Ⅲ. 평가개요 - 3. 평가 방법", ensure_report_text(blueprint.get("evaluationMethods"), "평가 방법")),
        ("eval-limitations", "Ⅲ. 평가개요 - 4. 평가의 한계", ensure_report_text(blueprint.get("evaluationLimitations"), "평가 한계")),
        ("eval-team", "Ⅲ. 평가개요 - 5. 평가단 구성 및 수행체계", ensure_report_text(blueprint.get("evaluationTeam"), "평가단 구성")),
        ("achievement", "Ⅳ. 성과달성도", ensure_report_text(blueprint.get("achievement"), "성과달성도")),
        ("criteria", "Ⅴ. 기준별 평가결과", criteria_findings_text(context, blueprint)),
        ("conclusion", "Ⅵ. 결론 - 1. 결론", ensure_report_text(blueprint.get("conclusion"), "결론")),
        (
            "factors",
            "Ⅵ. 결론 - 2. 작동요인 및 비작동요인",
            f"[작동요인]\n{list_text(blueprint.get('workingFactors'), '작동요인')}\n\n[비작동요인]\n{list_text(blueprint.get('nonWorkingFactors'), '비작동요인')}",
        ),
        ("theory", "Ⅵ. 결론 - 3. 변화이론 분석", ensure_report_text(blueprint.get("theoryOfChange"), "변화이론 분석")),
        ("feedback", "Ⅵ. 결론 - 4. 환류과제", feedback_tasks_text(blueprint.get("feedbackTasks"))),
        ("lessons", "Ⅵ. 결론 - 5. 교훈", lessons_text(blueprint.get("lessons"))),
        ("annex-en", "첨부 1. 평가결과 영문 요약", ensure_report_text(blueprint.get("englishSummary"), "영문 요약")),
        ("annex-field", "첨부 2. 현지(원격)조사 개요", ensure_report_text(blueprint.get("fieldSurveyOverview"), "현지 또는 원격 조사 개요")),
        ("annex-log", "첨부 3. 일별 활동내역", ensure_report_text(blueprint.get("activityLog"), "일별 활동내역")),
        ("annex-interview", "첨부 4. 면담자 목록 및 주요 면담 질문", ensure_report_text(blueprint.get("interviewQuestions"), "면담자 목록 및 질문")),
        ("annex-survey", "첨부 5. 설문조사지 및 설문조사 결과", ensure_report_text(blueprint.get("surveyResults"), "설문조사 결과")),
        ("annex-references", "첨부 6. 자료목록", reference_list_text(context)),
        ("annex-extra", "첨부 7. 그 외 첨부자료", ensure_report_text(blueprint.get("attachments"), "그 외 첨부자료")),
    ]
    return [
        {"id": section_id, "title": section_title, "body": saved.get(section_id) or body}
        for section_id, section_title, body in sections
    ]


def build_report_hwpx(context: dict, blueprint: dict, sections: list[dict] | None = None) -> bytes:
    """Build an editable HWPX report from the parsed HWP template structure."""
    paragraphs: list[tuple[str, int, int]] = []

    def add(text: str = "", style_id: int = 0, char_id: int = 0) -> None:
        for line in str(text or "").splitlines() or [""]:
            paragraphs.append((line, style_id, char_id))

    report_sections = sections or complete_report_sections(context, blueprint)
    for index, section in enumerate(report_sections):
        title = str(section.get("title", "")).strip()
        body_text = str(section.get("body", "")).strip()
        if index == 0:
            add(body_text or title, 1, 1)
        else:
            add(title, 1, 1)
            add(body_text or additional_info_note(title))
        add()

    body = "\n".join(hwpx_paragraph(text, 1000000000 + index, style_id, char_id) for index, (text, style_id, char_id) in enumerate(paragraphs))
    section_xml = f'''<?xml version="1.0" encoding="UTF-8"?>
<hs:sec xmlns:hs="http://www.hancom.co.kr/hwpml/2011/section"
        xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph">
<hp:secPr id="0" textDirection="HORIZONTAL" spaceColumns="0" tabStop="8000" tabStopVal="LEFT" tabStopUnit="HWPUNIT" outlineShapeIDRef="0" memoShapeIDRef="0" textVerticalWidthHead="0" masterPageCnt="0">
  <hp:grid lineGrid="0" charGrid="0" wonggojiFormat="0"/>
  <hp:startNum pageStartsOn="BOTH" page="0" pic="0" tbl="0" equation="0"/>
  <hp:visibility hideFirstHeader="0" hideFirstFooter="0" hideFirstMasterPage="0" border="SHOW_ALL" fill="SHOW_ALL" hideFirstPageNum="0" hideFirstEmptyLine="0" showLineNumber="0"/>
  <hp:lineNumberShape restartType="0" countBy="0" distance="0" startNumber="1"/>
  <hp:pagePr landscape="0" width="59528" height="84188" gutterType="LEFT_ONLY">
    <hp:margin header="2835" footer="2835" gutter="0" left="5669" right="5669" top="5669" bottom="5669"/>
  </hp:pagePr>
  <hp:footNotePr autoNumFormatType="DIGIT" autoNumFormatUserChar="" beforeDecorativeLetter="" afterDecorativeLetter=")" startNumber="1" numberingType="CONTINUOUS" placement="EACH_COLUMN" beneathText="0"/>
  <hp:endNotePr autoNumFormatType="DIGIT" autoNumFormatUserChar="" beforeDecorativeLetter="" afterDecorativeLetter=")" startNumber="1" numberingType="CONTINUOUS" placement="END_OF_DOCUMENT" beneathText="0"/>
  <hp:pageBorderFill type="BOTH" borderFillIDRef="0" textBorder="PAPER" headerInside="0" footerInside="0" fillArea="PAPER"/>
</hp:secPr>
{body}
</hs:sec>'''
    header_xml = '''<?xml version="1.0" encoding="UTF-8"?>
<hh:head xmlns:hh="http://www.hancom.co.kr/hwpml/2011/head">
  <hh:beginNum page="1" footnote="1" endnote="1" pic="1" tbl="1" equation="1"/>
  <hh:refList>
    <hh:fontfaces itemCnt="1"><hh:fontface lang="HANGUL" fontCnt="1"><hh:font id="0" type="ttf" name="함초롬바탕"/></hh:fontface></hh:fontfaces>
    <hh:charProperties itemCnt="2">
      <hh:charPr id="0" height="1000" textColor="#000000"/>
      <hh:charPr id="1" height="1300" textColor="#000000" bold="1"/>
    </hh:charProperties>
    <hh:paraProperties itemCnt="3">
      <hh:paraPr id="0" align="JUSTIFY"/>
      <hh:paraPr id="1" align="CENTER"/>
      <hh:paraPr id="2" align="LEFT"/>
    </hh:paraProperties>
    <hh:styles itemCnt="1"><hh:style id="0" type="PARA" name="Normal" engName="Normal" paraPrIDRef="0" charPrIDRef="0"/></hh:styles>
  </hh:refList>
</hh:head>'''
    content_hpf = '''<?xml version="1.0" encoding="UTF-8"?>
<opf:package xmlns:opf="http://www.idpf.org/2007/opf" version="3.0">
  <opf:metadata><opf:title>종료평가 결과보고서</opf:title></opf:metadata>
  <opf:manifest>
    <opf:item id="header" href="Contents/header.xml" media-type="application/xml"/>
    <opf:item id="section0" href="Contents/section0.xml" media-type="application/xml"/>
  </opf:manifest>
  <opf:spine><opf:itemref idref="section0"/></opf:spine>
</opf:package>'''
    version_xml = '''<?xml version="1.0" encoding="UTF-8"?><version app="ODA ImpactOps" version="1.0"/>'''
    container_xml = '''<?xml version="1.0" encoding="UTF-8"?><container xmlns="urn:oasis:names:tc:opendocument:xmlns:container" version="1.0"><rootfiles><rootfile full-path="Contents/content.hpf" media-type="application/hwpml-package+xml"/></rootfiles></container>'''
    output = BytesIO()
    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as hwpx:
        mime_info = zipfile.ZipInfo("mimetype")
        mime_info.compress_type = zipfile.ZIP_STORED
        hwpx.writestr(mime_info, "application/hwp+zip")
        hwpx.writestr("version.xml", version_xml)
        hwpx.writestr("META-INF/container.xml", container_xml)
        hwpx.writestr("Contents/content.hpf", content_hpf)
        hwpx.writestr("Contents/header.xml", header_xml)
        hwpx.writestr("Contents/section0.xml", section_xml)
        hwpx.writestr("Preview/PrvText.txt", "\n".join(text for text, _style, _char in paragraphs))
    return output.getvalue()


def fixed_width_text(value: str, width: int) -> str:
    text = str(value or "")
    if len(text) >= width:
        text = text[:width].rstrip(" ·-_")
        return text + (" " * (width - len(text)))
    return text + (" " * (width - len(text)))


def replace_hwp_body_text(template_path: Path, replacements: dict[str, str]) -> bytes:
    try:
        import olefile
    except Exception:
        raw = template_path.read_bytes()
        for source, target in replacements.items():
            raw = raw.replace(source.encode("utf-16le"), fixed_width_text(target, len(source)).encode("utf-16le"))
        return raw

    with tempfile.TemporaryDirectory(prefix="hwp_title_") as tmp:
        output_path = Path(tmp) / "report.hwp"
        shutil.copyfile(template_path, output_path)
        ole = olefile.OleFileIO(output_path, write_mode=True)
        try:
            section_paths = [
                entry for entry in ole.listdir()
                if len(entry) == 2 and entry[0] == "BodyText" and re.fullmatch(r"Section\d+", entry[1])
            ]
            section_paths.sort(key=lambda entry: int(entry[1].removeprefix("Section")))
            for stream_path in section_paths:
                sid = ole._find(stream_path)
                entry = ole.direntries[sid]
                raw_stream = ole.openstream(stream_path).read()
                try:
                    decompressed = zlib.decompress(raw_stream, -15)
                    wbits = -15
                except zlib.error:
                    decompressed = zlib.decompress(raw_stream)
                    wbits = zlib.MAX_WBITS
                updated = decompressed
                for source, target in replacements.items():
                    source_bytes = source.encode("utf-16le")
                    target_bytes = fixed_width_text(target, len(source)).encode("utf-16le")
                    updated = updated.replace(source_bytes, target_bytes)
                if updated == decompressed:
                    continue
                compressor = zlib.compressobj(level=6, wbits=wbits)
                packed = compressor.compress(updated) + compressor.flush()
                if len(packed) > entry.size:
                    continue
                packed = packed + (b"\x00" * (entry.size - len(packed)))
                if entry.size < ole.minisectorcutoff:
                    ole._write_mini_stream(entry, packed)
                else:
                    entry.build_sect_chain(ole)
                    position = 0
                    for sector in entry.sect_chain:
                        offset = (sector + 1) * ole.sectorsize
                        chunk = packed[position:position + ole.sectorsize]
                        ole.fp.seek(offset)
                        ole.fp.write(chunk.ljust(ole.sectorsize, b"\x00"))
                        position += ole.sectorsize
        finally:
            ole.close()

        raw = output_path.read_bytes()
        for source, target in replacements.items():
            raw = raw.replace(source.encode("utf-16le"), fixed_width_text(target, len(source)).encode("utf-16le"))
        return raw


def build_titled_hwp_report(context: dict) -> bytes:
    if not SAMPLE_REPORT_HWP_PATH.exists():
        return build_report_hwpx(context, fallback_report_blueprint(context))

    saved = read_report_editor_state()
    title_section = next((section for section in (saved or {}).get("sections", []) if section.get("id") == "title"), None)
    edited_title = (title_section or {}).get("body", "").splitlines()[0].replace(" 종료평가 결과보고서", "").strip() if title_section else ""
    project_title = edited_title or context["project"].get("title", "ODA 사업")
    report_title = fixed_width_text(f"{project_title} 종료평가 결과보고서", len("ㅇㅇ사업 종료평가 결과보고서"))
    cover_project = fixed_width_text(project_title, len("사업명(사업기간/예산)"))
    date_text = fixed_width_text(datetime.now().strftime("%Y.%m"), len("2023. 12"))
    return replace_hwp_body_text(SAMPLE_REPORT_HWP_PATH, {
        "ㅇㅇ사업 종료평가 결과보고서": report_title,
        "사업명(사업기간/예산)": cover_project,
        "2023. 12": date_text,
    })


def build_complete_hwpx_report(context: dict, blueprint: dict) -> bytes:
    return build_report_hwpx(context, blueprint, complete_report_sections(context, blueprint))


def build_report_plain_text(project: dict, criteria: list[dict], references: list[dict], overall: dict) -> str:
    lines = [
        "종료평가 결과보고서",
        "",
        f"사업명: {project.get('title', '사업명 확인 필요')}",
        f"사업기간: {project.get('period', '기간 확인 필요')}",
        f"사업비: {project.get('budget', '사업비 확인 필요')}",
        f"종합점수: {overall['score']}/{overall['maxScore']}점",
        f"KOICA 평가등급: {overall['koicaGrade']}",
        f"국무조정실 평가등급: {overall['governmentGrade']}",
        f"작성일: {now_label()}",
        "",
        "1. 평가결과",
        "",
    ]
    for index, criterion in enumerate(criteria, start=1):
        evaluation = criterion.get("evaluationResult") or {}
        criterion_id = criterion.get("id", "")
        criterion_refs = references_for_criterion(criterion_id, references)
        citation = " ".join(f"[{document['referenceNumber']}]" for document in criterion_refs[:5])
        evidence_line = f"주요 확인자료 {citation}".strip() if citation else "주요 확인자료는 자료목록 참조"
        lines += [
            f"{index}. {criterion_label(criterion)}({CRITERION_ENGLISH.get(criterion_id, '')}) 평가결과",
            evidence_line,
            "",
        ]
        lines.extend(f"- {line}" for line in plain_lines(evaluation.get("summary")))
        lines.append("")
        for section in evaluation.get("sections", []) or []:
            if section.get("title"):
                lines.append(section["title"])
            lines.extend(f"- {line}" for line in plain_lines(section.get("body")))
            lines.append("")
    lines += ["2. 자료목록", ""]
    if references:
        lines.extend(
            f"[{document['referenceNumber']}] {document.get('fileName', '')} | {document.get('criterionName', '')} | {document.get('evidenceName', '')}"
            for document in references
        )
    else:
        lines.append("등록된 자료 없음")
    return "\n".join(lines).strip() + "\n"


def build_template_hwp_report(project: dict, criteria: list[dict], references: list[dict], overall: dict) -> tuple[bytes | None, str | None]:
    if not HWP_REPORT_TEMPLATE_PATH.exists() or not HWP_REPORT_SCRIPT.exists():
        return None, "HWP template or automation script not found."
    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    content_path = REPORT_DIR / f"hwp_report_content_{timestamp}.txt"
    output_path = REPORT_DIR / f"원본양식_기반_종료평가_보고서_{timestamp}.hwp"
    content_path.write_text(build_report_plain_text(project, criteria, references, overall), encoding="utf-8")
    command = [
        "powershell.exe",
        "-NoProfile",
        "-ExecutionPolicy",
        "Bypass",
        "-File",
        str(HWP_REPORT_SCRIPT),
        "-TemplatePath",
        str(HWP_REPORT_TEMPLATE_PATH),
        "-OutputPath",
        str(output_path),
        "-ContentPath",
        str(content_path),
    ]
    try:
        subprocess.run(command, check=True, capture_output=True, text=True, timeout=90)
    except Exception as exc:  # noqa: BLE001
        return None, f"HWP generation failed: {exc}"
    if not output_path.exists() or output_path.stat().st_size <= 0:
        return None, "HWP generation completed without an output file."
    return output_path.read_bytes(), None


def build_evaluation_report_package() -> tuple[bytes, str]:
    attach_uploaded_documents()
    apply_persisted_evaluations()
    project = project_payload()
    criteria = report_criteria()
    references = reference_documents_for_report()
    total_score = round(sum(float(item.get("currentScore4", 1) or 1) for item in criteria), 1)
    koica_grade, government_grade = grade_label(total_score)
    overall = {
        "score": total_score,
        "maxScore": 20,
        "koicaGrade": koica_grade,
        "governmentGrade": government_grade,
    }
    context = build_report_context(project, criteria, references, overall)
    blueprint = generate_report_blueprint(context)
    xlsx_bytes = build_template_grade_xlsx(context, blueprint)
    pptx_bytes = build_template_lesson_pptx(context, blueprint)
    from ..hwpx.patchers import build_cover_grade_body_patched_hwpx
    patched_report = build_cover_grade_body_patched_hwpx()
    report_bytes = base64.b64decode(patched_report["data"])
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    project_name = safe_filename_part(project.get("title", "ODA_사업"))
    package_name = f"{project_name}_종료평가_보고서_패키지_{timestamp}.zip"
    output = BytesIO()
    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as package:
        package.writestr(f"{project_name}_5-1_종료평가_결과보고서.hwpx", report_bytes)
        package.writestr(f"{project_name}_5-2_종료평가_등급_결과표.xlsx", xlsx_bytes)
        if pptx_bytes:
            package.writestr(f"{project_name}_5-3_분야별_평가_교훈_리포트.pptx", pptx_bytes)
    return output.getvalue(), package_name


def current_report_context() -> dict:
    attach_uploaded_documents()
    apply_persisted_evaluations()
    project = project_payload()
    criteria = report_criteria()
    references = reference_documents_for_report()
    total_score = round(sum(float(item.get("currentScore4", 1) or 1) for item in criteria), 1)
    koica_grade, government_grade = grade_label(total_score)
    overall = {
        "score": total_score,
        "maxScore": 20,
        "koicaGrade": koica_grade,
        "governmentGrade": government_grade,
    }
    return build_report_context(project, criteria, references, overall)


def current_report_context_and_blueprint() -> tuple[dict, dict]:
    context = current_report_context()
    return context, generate_report_blueprint(context)

