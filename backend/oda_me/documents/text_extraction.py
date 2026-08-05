from __future__ import annotations

from ..core import *


def file_signature(path_value: str | None) -> dict:
    if not path_value:
        return {}
    path = Path(path_value)
    try:
        stat = path.stat()
    except OSError:
        return {"name": path.name, "missing": True}
    return {
        "name": path.name,
        "size": stat.st_size,
        "mtime": round(stat.st_mtime, 3),
    }


def normalize_extracted_text(text: str) -> str:
    readable = "".join(char if char.isprintable() or char in "\n\r\t" else " " for char in text)
    readable = re.sub(r"\r\n?", "\n", readable)
    readable = re.sub(r"[ \t]+", " ", readable)
    readable = re.sub(r"\n{3,}", "\n\n", readable)
    return readable.strip()


def collect_rhwp_strings(value: object) -> list[str]:
    if isinstance(value, str):
        return [value]
    if isinstance(value, list):
        collected: list[str] = []
        for item in value:
            collected.extend(collect_rhwp_strings(item))
        return collected
    if isinstance(value, dict):
        collected = []
        for item in value.values():
            collected.extend(collect_rhwp_strings(item))
        return collected
    return []


def meaningful_rhwp_fragment(value: str) -> bool:
    cleaned = value.strip()
    if len(cleaned) < 2:
        return False
    if re.fullmatch(r"[-+]?[\d.,:/\\ ]+", cleaned):
        return False
    if re.search(r"[\uac00-\ud7a3]", cleaned):
        return True
    return len(cleaned) >= 12 and bool(re.search(r"\s", cleaned))


def extract_text_from_rhwp_output(output: str) -> str:
    normalized = normalize_extracted_text(output)
    if not normalized:
        return ""
    try:
        parsed = json.loads(normalized)
    except json.JSONDecodeError:
        quoted = re.findall(r'"((?:[^"\\]|\\.)*)"', normalized)
        fragments: list[str] = []
        for item in quoted:
            try:
                item = bytes(item, "utf-8").decode("unicode_escape")
            except UnicodeError:
                pass
            if meaningful_rhwp_fragment(item):
                fragments.append(item)
        if fragments:
            return normalize_extracted_text("\n".join(fragments))
        lines = [line.strip() for line in normalized.splitlines() if meaningful_rhwp_fragment(line)]
        return normalize_extracted_text("\n".join(lines))
    fragments = [item for item in collect_rhwp_strings(parsed) if meaningful_rhwp_fragment(item)]
    return normalize_extracted_text("\n".join(fragments))


def extract_hwp_text_with_rhwp(raw: bytes, filename: str) -> tuple[str, str] | None:
    suffix = Path(filename).suffix.lower()
    if suffix not in HWP_SUFFIXES:
        return None
    with tempfile.TemporaryDirectory(prefix="rhwp_") as tmp:
        tmp_path = Path(tmp)
        input_path = Path(tmp) / f"input{suffix}"
        output_dir = tmp_path / "text"
        input_path.write_bytes(raw)
        commands = (
            ("export-text", str(input_path), "-o", str(output_dir)),
            ("export-markdown", str(input_path), "-o", str(output_dir)),
            ("dump", str(input_path)),
        )
        last_error = ""
        for command in commands:
            try:
                result = subprocess.run(
                    [RHWP_BIN, *command],
                    capture_output=True,
                    text=True,
                    encoding="utf-8",
                    errors="replace",
                    timeout=RHWP_TIMEOUT_SECONDS,
                    check=False,
                )
            except (FileNotFoundError, subprocess.SubprocessError, OSError) as exc:
                last_error = str(exc)
                break
            if result.returncode != 0:
                last_error = result.stderr.strip() or result.stdout.strip()
                continue
            exported_files = sorted(output_dir.glob("*.txt")) + sorted(output_dir.glob("*.md"))
            if exported_files:
                exported_text = "\n\n".join(file.read_text(encoding="utf-8", errors="replace") for file in exported_files)
                extracted = normalize_extracted_text(exported_text)
                if extracted:
                    return extracted[:30000], f"rhwp_{command[0]}"
            extracted = extract_text_from_rhwp_output(result.stdout)
            if extracted:
                return extracted[:30000], f"rhwp_{command[0]}"
            output_dir.mkdir(exist_ok=True)
            for file in output_dir.iterdir():
                if file.is_file():
                    file.unlink()
        if last_error:
            return f"rhwp HWP/HWPX parser failed: {last_error}", "rhwp_failed"
    return None


def extract_text(raw: bytes, filename: str, mime_type: str) -> tuple[str, str]:
    suffix = Path(filename).suffix.lower()
    rhwp_result = extract_hwp_text_with_rhwp(raw, filename)
    if rhwp_result and rhwp_result[1] != "rhwp_failed":
        return rhwp_result
    if suffix in {".txt", ".md", ".csv", ".json", ".html", ".xml"} or mime_type.startswith("text/"):
        for encoding in ("utf-8-sig", "utf-8", "cp949", "euc-kr", "utf-16"):
            try:
                return normalize_extracted_text(raw.decode(encoding)), "decoded_text"
            except UnicodeDecodeError:
                continue
    utf16_text = normalize_extracted_text(raw.decode("utf-16le", errors="ignore"))
    if any(token in utf16_text for token in ("사업개요서", "사업규모", "사업기간", "사업비", "사업명")):
        return utf16_text[:30000], "utf16le_strings"
    decoded = raw.decode("latin-1", errors="ignore")
    fragments = re.findall(r"[A-Za-z0-9가-힣 .,;:()/%+\-_\[\]\n\r]{8,}", decoded)
    text = "\n".join(fragment.strip() for fragment in fragments if fragment.strip())
    if text:
        return normalize_extracted_text(text)[:20000], "binary_printable_strings"
    if rhwp_result:
        return rhwp_result
    return "텍스트 자동 추출이 필요한 바이너리 문서입니다. PDF/HWP/DOCX 전용 파서 연동 후 본문 추출 품질을 개선할 수 있습니다.", "needs_parser"


def cache_key_for_document(document: dict) -> str:
    raw_signature = file_signature(document.get("rawPath"))
    payload = {
        "id": document.get("id"),
        "fileName": document.get("fileName"),
        "size": raw_signature.get("size") or document.get("size"),
        "mtime": raw_signature.get("mtime"),
        "method": "full_reference_v2",
    }
    return hashlib.sha256(json.dumps(payload, ensure_ascii=False, sort_keys=True).encode("utf-8")).hexdigest()


def extract_pdf_text_full(path: Path) -> tuple[str, str]:
    try:
        from pypdf import PdfReader
    except Exception as exc:  # noqa: BLE001
        return "", f"pdf_parser_unavailable:{exc}"
    try:
        reader = PdfReader(str(path))
        pages = []
        for index, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"[page {index}]\n{text}")
        return normalize_extracted_text("\n\n".join(pages)), "pypdf_extract_text"
    except Exception as exc:  # noqa: BLE001
        return "", f"pypdf_failed:{exc}"


def extract_xlsx_text_full(path: Path) -> tuple[str, str]:
    try:
        from openpyxl import load_workbook
    except Exception as exc:  # noqa: BLE001
        return "", f"xlsx_parser_unavailable:{exc}"
    try:
        workbook = load_workbook(str(path), data_only=True, read_only=True)
        chunks = []
        for sheet in workbook.worksheets:
            chunks.append(f"[sheet] {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if values:
                    chunks.append(" | ".join(values))
        workbook.close()
        return normalize_extracted_text("\n".join(chunks)), "openpyxl_values"
    except Exception as exc:  # noqa: BLE001
        return "", f"openpyxl_failed:{exc}"


def extract_pptx_text_full(path: Path) -> tuple[str, str]:
    try:
        from pptx import Presentation
    except Exception as exc:  # noqa: BLE001
        return "", f"pptx_parser_unavailable:{exc}"
    try:
        presentation = Presentation(str(path))
        chunks = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            chunks.append(f"[slide {slide_index}]")
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    chunks.append(text)
        return normalize_extracted_text("\n".join(chunks)), "python_pptx_text"
    except Exception as exc:  # noqa: BLE001
        return "", f"python_pptx_failed:{exc}"


def extract_hwp_reference_text_full(path: Path) -> tuple[str, str]:
    suffix = path.suffix.lower()
    if suffix not in HWP_SUFFIXES:
        return "", "not_hwp"
    with tempfile.TemporaryDirectory(prefix="rhwp_full_") as tmp:
        output_dir = Path(tmp) / "text"
        commands = (
            ("export-text", str(path), "-o", str(output_dir)),
            ("export-markdown", str(path), "-o", str(output_dir)),
            ("dump", str(path)),
        )
        last_error = ""
        for command in commands:
            try:
                result = subprocess.run(
                    [RHWP_BIN, *command],
                    capture_output=True,
                    text=True,
                    encoding="utf-8",
                    errors="replace",
                    timeout=RHWP_TIMEOUT_SECONDS,
                    check=False,
                )
            except (FileNotFoundError, subprocess.SubprocessError, OSError) as exc:
                return "", f"rhwp_full_failed:{exc}"
            if result.returncode != 0:
                last_error = result.stderr.strip() or result.stdout.strip()
                continue
            exported_files = sorted(output_dir.glob("*.txt")) + sorted(output_dir.glob("*.md"))
            if exported_files:
                exported_text = "\n\n".join(file.read_text(encoding="utf-8", errors="replace") for file in exported_files)
                extracted = normalize_extracted_text(exported_text)
                if extracted:
                    return extracted, f"rhwp_full_{command[0]}"
            extracted = extract_text_from_rhwp_output(result.stdout)
            if extracted:
                return extracted, f"rhwp_full_{command[0]}"
        return "", f"rhwp_full_failed:{last_error}"


def read_reference_text_full(document: dict) -> dict:
    raw_path_value = document.get("rawPath")
    text_path_value = document.get("textPath")
    cache_key = cache_key_for_document(document)
    cache_path = REFERENCE_TEXT_CACHE_DIR / f"{cache_key}.json"
    if cache_path.exists():
        try:
            return json.loads(cache_path.read_text(encoding="utf-8-sig"))
        except Exception:
            pass

    raw_path = Path(raw_path_value) if raw_path_value else None
    suffix = Path(str(document.get("fileName") or raw_path_value or "")).suffix.lower()
    text = ""
    method = "missing_source"
    if raw_path and raw_path.exists():
        if suffix == ".pdf":
            text, method = extract_pdf_text_full(raw_path)
        elif suffix in {".xlsx", ".xlsm"}:
            text, method = extract_xlsx_text_full(raw_path)
        elif suffix == ".pptx":
            text, method = extract_pptx_text_full(raw_path)
        elif suffix in HWP_SUFFIXES:
            text, method = extract_hwp_reference_text_full(raw_path)
        elif suffix in {".txt", ".md", ".csv", ".json", ".html", ".xml"}:
            text = raw_path.read_text(encoding="utf-8", errors="replace")
            method = "raw_text_file"

    if not text and text_path_value and Path(text_path_value).exists():
        text = Path(text_path_value).read_text(encoding="utf-8", errors="replace")
        method = f"cached_{document.get('extractionMethod') or 'textPath'}"
    text = normalize_extracted_text(text)
    korean_chars = len(re.findall(r"[\uac00-\ud7a3]", text))
    alpha_chars = len(re.findall(r"[A-Za-z]", text))
    if not text:
        quality = "empty"
    elif method.startswith("pypdf") and len(text) < 500:
        quality = "ocr_required"
    elif "binary_printable" in method or text.startswith("%PDF") or "PK" in text[:500]:
        quality = "low"
    elif korean_chars + alpha_chars < 500:
        quality = "low"
    else:
        quality = "usable"
    payload = {
        "text": text,
        "method": method,
        "quality": quality,
        "charCount": len(text),
        "koreanChars": korean_chars,
        "alphaChars": alpha_chars,
    }
    REFERENCE_TEXT_CACHE_DIR.mkdir(parents=True, exist_ok=True)
    cache_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return payload


def truncate_reference_text(text: str, limit: int) -> tuple[str, bool]:
    text = normalize_extracted_text(text)
    if len(text) <= limit:
        return text, False
    if limit <= 0:
        return "", True
    if limit < 2000:
        return text[:limit].rstrip(), True
    head_limit = max(1000, int(limit * 0.72))
    tail_limit = max(800, limit - head_limit)
    truncated = (
        text[:head_limit].rstrip()
        + "\n\n[...문서 중간 일부 생략: 컨텍스트 한도 내에서 앞부분과 끝부분을 우선 제공...]\n\n"
        + text[-tail_limit:].lstrip()
    )
    if len(truncated) > limit:
        truncated = truncated[:limit].rstrip()
    return truncated, True

